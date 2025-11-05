
"""
生成 attention_cem_briefing.pptx
脚本聚焦清晰、精简、视觉友好的呈现。
运行前请确保已安装 python-pptx。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

# -------------------------------------------------------------------
# 样式参数
# -------------------------------------------------------------------
BG_PRIMARY = RGBColor(16, 44, 87)      # 深蓝封面
BG_LIGHT = RGBColor(242, 245, 250)     # 浅灰背景
ACCENT_GREEN = RGBColor(0, 158, 115)   # 亮绿色强调
ACCENT_BLUE = RGBColor(33, 150, 243)   # 蓝色强调
ACCENT_ORANGE = RGBColor(255, 153, 0)  # 对比色
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_with_line(slide, title_text, subtitle_text=None, accent=ACCENT_GREEN):
    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(20, 20, 20)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.8), Inches(1.8), Inches(0.35), Inches(0.35)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.color.rgb = accent

    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(8.0), Inches(1.0))
        p = subtitle_box.text_frame.add_paragraph()
        p.text = subtitle_text
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(80, 80, 80)

def add_bullets(slide, bullets, left=1.0, top=2.4, width=8.0, height=3.5, font_size=24):
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.word_wrap = True
    for idx, (text, highlight) in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = BODY_FONT
        if highlight:
            p.font.color.rgb = highlight
        if idx == 0:
            tf.paragraphs[0].text = text
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.name = BODY_FONT
            if highlight:
                tf.paragraphs[0].font.color.rgb = highlight

def add_note(slide, note_text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = note_text

def add_process_diagram(slide, boxes, top=2.3):
    """
    boxes: list of tuples (text, fill_color)
    """
    left = 1.0
    spacing = 0.4
    for text, color in boxes:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(2.6), Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.text = text
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.name = BODY_FONT
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        left += 2.6 + spacing

        if left > 8.0:
            top += 1.6
            left = 1.0

def build_presentation():
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.333)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_PRIMARY)
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.5), Inches(2.5))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = "Attention-CEM 防御框架概览"
    p.font.size = Pt(54)
    p.font.name = TITLE_FONT
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.5), Inches(1.0))
    p2 = sub_box.text_frame.paragraphs[0]
    p2.text = "Slot + Gated Cross Attention  vs.  Gated Attention Pooling"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(230, 230, 230)
    add_note(slide, "开场说明研究背景与两套注意力替换方案的目标。")

    # 议程
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, BG_LIGHT)
    add_title_with_line(slide, "会议导览", "把握结构，突出两套防御的关键差异", ACCENT_BLUE)
    add_bullets(slide, [
        ("协同推理隐私问题 & CEM 目标", None),
        ("Slot + Gated Cross Attention：结构与训练融合", None),
        ("Gated Attention Pooling：轻量化替代方案", None),
        ("对比结论 & 下一步验证计划", None),
    ])
    add_note(slide, "说明将从威胁模型与CEM需求讲起，再分别解读两套框架，最后给出决策指南。")

    # 背景
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, RGBColor(255, 255, 255))
    add_title_with_line(slide, "协同推理威胁与CEM目标", accent=ACCENT_GREEN)
    add_bullets(slide, [
        ("Client 侧编码器 Fe → 中间特征 z → Cloud 侧 Fd", None),
        ("攻击者利用 z 做模型反演；关键指标是 H(x | z)", None),
        ("原方案：KMeans/GMM 近似条件熵，难捕捉类内多模态", ACCENT_ORANGE),
    ], top=2.3)
    add_bullets(slide, [
        ("新策略：以注意力模块替换近似器", ACCENT_GREEN),
        ("优势：可训练、可微分、兼容现有优化流程", None),
    ], top=4.4)
    add_note(slide, "引出注意力替换的原因：需要更丰富的类内建模方式。")

    # Slot + Gated Cross Attention 总览
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, RGBColor(255, 255, 255))
    add_title_with_line(slide, "Slot + Gated Cross Attention 框架鸟瞰", accent=ACCENT_BLUE)
    add_process_diagram(slide, [
        ("类内特征 Tokens", ACCENT_BLUE),
        ("Slot Attention\n多轮竞争", ACCENT_GREEN),
        ("Gated Cross Attention\n注入 slot 语境", ACCENT_ORANGE),
        ("门控方差堆栈\n输出 CEM loss", RGBColor(157, 114, 255)),
    ])
    add_bullets(slide, [
        ("每个类独立处理，slots 复用以提取潜在模态", None),
        ("Flamingo 风格门控残差，控制 cross-attn/FFN 渐进启用", None),
        ("最终只回传类内 log-variance，嵌入现有 CEM 流程", None),
    ], top=4.5, font_size=22)
    add_note(slide, "强调图中四步流程，指出 slots 如何作为记忆池提升类内建模。")

    # Slot Attention 内部细节
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, RGBColor(255, 255, 255))
    add_title_with_line(slide, "Slot Attention：如何抽取类内子模态", accent=ACCENT_GREEN)
    add_bullets(slide, [
        ("LayerNorm + 一次性 K/V 投影，减少重复开销", None),
        ("学习到的 (μ, σ) 初始化 + 3 轮 GRU 更新强化 slot", None),
        ("softmax(β · sim) 竞争调整 slot 覆盖度，β 可学习", None),
        ("残差 MLP 稳定 slot 分布，避免硬切换", None),
    ], top=2.4, font_size=23)
    add_note(slide, "讲重点：竞争式 softmax 和 β 学习，说明 Slot Attention 能捕捉多模态结构。")

    # Gated Cross Attention & 门控堆栈
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, RGBColor(255, 255, 255))
    add_title_with_line(slide, "Gated Cross Attention + 门控熵堆栈", accent=ACCENT_ORANGE)
    add_bullets(slide, [
        ("Pre-LN 多头交叉注意力，α_attn 控制残差权重", None),
        ("FFN 前置层归一化，α_ffn 渐进激活 dense 层", None),
        ("方差估计：slot 责任 → μ_s, σ_s^2 → 多级门控", None),
    ], top=2.4, font_size=23)
    add_bullets(slide, [
        ("Per-dim gate：sigmoid(MLP(LN(log σ_s^2)))", ACCENT_BLUE),
        ("SNR gate：σ(κ(σ_s^2/(μ_s^2+ε) − τ))", ACCENT_ORANGE),
        ("Softplus margin + slot mass 权重 + class gate", None),
    ], top=4.4, font_size=21)
    add_note(slide, "快速扫过三层门控：维度gate、SNR gate、类gate，突出控制策略。")

    # Slot 框架训练整合
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_with_line(slide, "Slot 框架：训练流程 & 安全阀", accent=ACCENT_BLUE)
    add_bullets(slide, [
        ("激活条件：稳定后再开启（warmup > 3 epoch, λ > 0）", None),
        ("梯度路径：rob loss 反传 → 缓存 encoder grad → 再合并", None),
        ("缺省 scaling：attention_loss_scale = 0.25", None),
    ], top=2.4)
    add_bullets(slide, [
        ("Early shutoff：首 100 次或 gate 超阈值 → 输出 0", ACCENT_ORANGE),
        ("异常处理：NaN/Inf 直接跳过，打印 gate 统计日志", None),
        ("优化器注册：首次调用时添加参数组，避免空算力", None),
    ], top=4.4)
    add_note(slide, "强调安全阀：Early shutoff + 日志，有助指导调参。")

    # Gated Attention Pooling 总览
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_with_line(slide, "Gated Attention Pooling 框架鸟瞰", accent=ACCENT_GREEN)
    add_process_diagram(slide, [
        ("类内特征 X_c\nLayerNorm", ACCENT_BLUE),
        ("Gated Pooling\n tanh(Vx) · σ(Ux)", ACCENT_GREEN),
        ("Softmax 权重\na_m", ACCENT_ORANGE),
        ("加权 μ, σ²\n→ variance hinge", RGBColor(120, 94, 240)),
    ])
    add_bullets(slide, [
        ("结构极简：两层线性 + softmax", None),
        ("Variance hinge：log(σ_c^2) 与阈值比较", None),
        ("定位：快速验证、大 batch 或算力紧张场景", None),
    ], top=4.5, font_size=22)
    add_note(slide, "突出简单流程和典型使用场景。")

    # Gated Pooling 数学 & 训练
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_with_line(slide, "Gated Pooling 数学细节", accent=ACCENT_BLUE)
    add_bullets(slide, [
        ("a_m ∝ exp(wᵀ[tanh(Vx_m) ⊙ σ(Ux_m)])", None),
        ("μ_c = Σ a_m x_m,  σ_c² = Σ a_m (x_m − μ_c)²", None),
        ("τ = var_threshold · reg_strength² + γ", None),
        ("L_C = max{0, log(σ_c² + γ) − log(τ)}", ACCENT_ORANGE),
    ], top=2.4, font_size=22)
    add_bullets(slide, [
        ("LayerNorm 防止注意力塌缩", None),
        ("默认 hidden_dim = clip(D/4)", None),
        ("attention_loss_scale = 0.1", None),
        ("skip 样本数 ≤1 的类避免噪声", None),
    ], top=4.4, font_size=22)
    add_note(slide, "强调阈值 τ 与 LayerNorm 作用。")

    # Gated Pooling 训练融入
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_with_line(slide, "轻量方案如何融入训练循环", accent=ACCENT_GREEN)
    add_bullets(slide, [
        ("Warmup 更长：默认 5 epoch 保障权重稳定", None),
        ("首次使用时自动实例化 + 注册参数组", None),
        ("NaN/Inf 处理与 Slot 框架一致，易维护", None),
    ], top=2.4)
    add_note(slide, "对比 slot 方案：保留关键守则，但计算复杂度更低。")

    # 对比总结
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_with_line(slide, "两种注意力 CEM 方案对比", accent=ACCENT_ORANGE)
    add_bullets(slide, [
        ("表达能力：Slot → 精细多模态；Pooling → 汇聚摘要", None),
        ("计算开销：Slot 重（GRU + 多头）；Pooling 轻", None),
        ("调参：Slot 多门控；Pooling 关注阈值与缩放", None),
        ("安全守则：Slot 需 shutoff；Pooling 简洁但须监控 softmax", None),
    ], top=2.4, font_size=23)
    add_note(slide, "点出决策依据，便于导师快速判断使用场景。")

    # 下一步计划
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title_with_line(slide, "下一步验证计划", accent=ACCENT_BLUE)
    add_bullets(slide, [
        ("多数据集对比：CIFAR-10/100, FaceScrub, TinyImageNet", None),
        ("与 GMM 基线比对 MIA 指标 (MSE/SSIM)", None),
        ("Mixed 策略：先 Pooling 筛选 → Slot 精细雕琢", None),
        ("分析 Early shutoff / class gate 对隐私-准确率折衷的影响", None),
    ], top=2.4, font_size=23)
    add_note(slide, "收尾强调评测计划与潜在混合策略。")

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_PRIMARY)
    thank_box = slide.shapes.add_textbox(Inches(3.0), Inches(2.5), Inches(7.0), Inches(2.0))
    tp = thank_box.text_frame.paragraphs[0]
    tp.text = "感谢指导 · 欢迎讨论"
    tp.font.size = Pt(40)
    tp.font.color.rgb = RGBColor(255, 255, 255)
    tp.font.bold = True
    tp.font.name = TITLE_FONT
    add_note(slide, "邀请反馈，视情况衔接至问答。")

    prs.save("attention_cem_briefing.pptx")
    print("已生成：attention_cem_briefing.pptx")

if __name__ == "__main__":
    build_presentation()